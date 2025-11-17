from langchain.docstore.document import Document
from pptx import Presentation
import os
from backend.utilis import *
from backend.Image_processing_disk import extract_Image_summaries
from unstructured.partition.pptx import partition_pptx
import re
import asyncio
from src.logger_config import log
import re
import tempfile
import io
import shutil


def pptx_processor(file_bytes):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(file_bytes)
            tmp.flush()                  
            tmp_path = tmp.name
        log.info(f"Processing temp PDF: {tmp_path}")

        raw_pptx_elements = partition_pptx(
        filename = tmp_path,
        include_page_breaks = True,
        infer_table_structure = True,
        starting_page_number = 1,
        strategy = "hi-res")
        return raw_pptx_elements
    
    except Exception:
        log.exception("processing ppt failed")
        return []
        
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        log.info(f"Tmp File removed: {tmp_path}")



def extract_images_from_pptx(file_bytes: str, output_dir):
    """Extract embedded images from DOCX and save them to output_dir."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        images = []
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    img = shape.image
                    images.append({
                        "slide": i+1,
                        "image_bytes": img.blob,
                        "ext": img.ext
                    })

        if not images:
            return []
        
        img_paths = []

        for idx, img_info in enumerate(images):
            slide_no = img_info["slide"]
            ext = img_info["ext"]
            img_bytes = img_info["image_bytes"]
            img_path = os.path.join(output_dir, f"slide_{slide_no}_{idx}.{ext}")
            with open(img_path, "wb") as f:
                f.write(img_bytes)
            img_paths.append({"path" : img_path})
        return img_paths    
    except Exception:
        log.exception("Image extraction failed")
        return []



def extract_pptx_elements(file_name,file_bytes,user_id):
    try:
        if not file_bytes:
            log.warning(f"[Excel_PARSE] Empty file: {file_name}")
            return [], "Uploaded excel is empty."
        
        log.info("Enter processor function to extract raw elements from pptx")
        raw_pptx_elements = pptx_processor(file_bytes)
        log.info(f"Raw_pdf_elements {raw_pptx_elements}")

        if not raw_pptx_elements:
            log.warning(f"No elements found in {file_name}")
            return [], "No readable text or elements found in the PPT."

        log.info("Get the folder to store extracted image from pptx")
        output_dir = get_doc_image_dir(file_name,user_id)
        log.info(f"Directory path {output_dir}")


        log.info("Enter processor function to extract text elements from raw elements")
        try:
         Header, Footer, Title , NarrativeText, Text , ListItem , Img , Tables = extract_text_elements(raw_pptx_elements)
        except Exception:
            log.exception("[PPT_PARSE] Text extraction failed.")
        
        pptx_image_info  = extract_images_from_pptx(file_bytes, output_dir)

        Image_summaries = []

        if pptx_image_info:
            log.info('Pass extracted images for summarization')
            Image_summaries = asyncio.run(extract_Image_summaries(output_dir))
            if Image_summaries:
                 Image_summaries = [re.sub(r'[>]+', '', t).strip() for t in Image_summaries if t.strip()]
                 log.info(f"cleaned image summaries extracted successfully. Count: {len(Image_summaries)}")
            else:
                Image_summaries = []

        final = final_doc(Header, Footer, Title , NarrativeText, Text , ListItem , Img , Tables, Image_summaries, file_name)
        documents = final.overall()

        if not documents:
            log.warning(f"PDF parsed but no valid text documents found in {file_name}")
            return [], "No readable text detected."

        log.info(f"Extracted {len(documents)} document chunks from {file_name}")
        return documents, None  # No error message
    
    except Exception as e:
        log.exception("Creating documents object failed")
        return [], f"PPT extraction failed: {e}"



